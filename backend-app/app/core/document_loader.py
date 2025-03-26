import logging
from langchain_community.document_loaders import PyPDFLoader
from langchain_community.document_loaders.csv_loader import CSVLoader
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders import JSONLoader
from typing import List, Dict, Any, Optional

# Use logging instead of fastapi.logger
logger = logging.getLogger(__name__)

def load_pdf_to_text(file_path: str) -> List[str]:
    """
    Load a PDF file and extract its text content.

    Parameters:
        file_path: The path to the PDF file.

    Returns:
        A list of strings, where each string is a page of text from the PDF.
    """
    try:
        # Initialize the PDF loader
        loader = PyPDFLoader(file_path)
        
        # Load and extract text from the PDF
        documents = loader.load()
        
        # Extract text from each page
        text_pages = [doc.page_content for doc in documents]
        
        return text_pages
    except Exception as e:
        logger.error(f"Error loading PDF file: {str(e)}")
        return []

def load_csv_to_text(file_path: str, csv_args: Optional[Dict[str, Any]] = None) -> List[str]:
    """
    Load a CSV file and extract its text content.

    Parameters:
        file_path: The path to the CSV file.
        csv_args: Optional arguments for CSV parsing (delimiter, quotechar, etc.)

    Returns:
        A list of strings, where each string represents a row from the CSV.
    """
    if csv_args is None:
        csv_args = {"delimiter": ",", "quotechar": '"'}
    
    try:
        # Initialize the CSV loader with optional arguments
        loader = CSVLoader(
            file_path=file_path,
            csv_args=csv_args
        )
        
        # Load and extract text from the CSV
        documents = loader.load()
        
        # Extract text from each row
        text_rows = [doc.page_content for doc in documents]
        
        return text_rows
    except Exception as e:
        logger.error(f"Error loading CSV file: {str(e)}")
        return []

def load_excel_to_text(file_path: str) -> List[str]:
    """
    Load an Excel file and extract its text content.

    Parameters:
        file_path: The path to the Excel file.

    Returns:
        A list of strings, where each string represents content from the Excel file.
    """
    try:
        # Try to import pandas first
        try:
            import pandas as pd
            
            # Read Excel file using pandas
            df = pd.read_excel(file_path)
            
            # Convert DataFrame to text
            text_content = []
            
            # Add column headers
            headers = df.columns.tolist()
            text_content.append(", ".join(str(header) for header in headers))
            
            # Add each row
            for _, row in df.iterrows():
                row_text = ", ".join(str(value) for value in row.values)
                text_content.append(row_text)
            
            return text_content
            
        except ImportError:
            # If pandas is not available, try using openpyxl directly
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path)
            text_content = []
            
            for sheet in workbook.worksheets:
                sheet_text = f"Sheet: {sheet.title}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ", ".join(str(cell) if cell is not None else "" for cell in row)
                    sheet_text += row_text + "\n"
                
                text_content.append(sheet_text)
            
            return text_content
            
    except Exception as e:
        logger.error(f"Error loading Excel file: {str(e)}")
        return []

def load_ppt_to_text(file_path: str) -> List[str]:
    """
    Load a PowerPoint file and extract its text content.

    Parameters:
        file_path: The path to the PowerPoint file.

    Returns:
        A list of strings, where each string represents content from a slide.
    """
    try:
        # Use python-pptx directly instead of UnstructuredPowerPointLoader
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = f"Slide {i+1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                text_content.append(slide_text)
            
            return text_content
        except ImportError:
            logger.warning("python-pptx not installed. Falling back to alternative method.")
            # You could implement a fallback method here if needed
            return ["PowerPoint extraction requires python-pptx library"]
    except Exception as e:
        logger.error(f"Error loading PowerPoint file: {str(e)}")
        return []

def load_docx_to_text(file_path: str) -> List[str]:
    """
    Load a Word document and extract its text content.

    Parameters:
        file_path: The path to the Word document.

    Returns:
        A list containing the text content of the Word document.
    """
    try:
        # Initialize the Word document loader
        loader = Docx2txtLoader(file_path)
        
        # Load and extract text from the Word document
        documents = loader.load()
        
        # Extract text from the document
        text_content = [doc.page_content for doc in documents]
        
        return text_content
    except Exception as e:
        logger.error(f"Error loading Word document: {str(e)}")
        return []

def load_document_to_text(file_path: str, file_type: Optional[str] = None, **kwargs) -> List[str]:
    """
    Load a document file and extract its text content based on file type.

    Parameters:
        file_path: The path to the document file.
        file_type: The type of the file ('pdf', 'csv', 'excel', 'ppt', 'docx', 'json'). 
                   If None, will be inferred from file extension.
        **kwargs: Additional arguments for specific loaders.

    Returns:
        A list of strings containing the extracted text.
    """
    if file_type is None:
        # Infer file type from extension
        if file_path.lower().endswith('.pdf'):
            file_type = 'pdf'
        elif file_path.lower().endswith('.csv'):
            file_type = 'csv'
        elif file_path.lower().endswith(('.xlsx', '.xls')):
            file_type = 'excel'
        elif file_path.lower().endswith(('.pptx', '.ppt')):
            file_type = 'ppt'
        elif file_path.lower().endswith(('.docx', '.doc')):
            file_type = 'docx'
        elif file_path.lower().endswith('.json'):
            file_type = 'json'
        else:
            logger.error(f"Unsupported file type for {file_path}")
            return []
    
    if file_type.lower() == 'pdf':
        return load_pdf_to_text(file_path)
    elif file_type.lower() == 'csv':
        csv_args = kwargs.get('csv_args', {})
        return load_csv_to_text(file_path, csv_args)
    elif file_type.lower() == 'excel':
        return load_excel_to_text(file_path)
    elif file_type.lower() == 'ppt':
        return load_ppt_to_text(file_path)
    elif file_type.lower() == 'docx':
        return load_docx_to_text(file_path)
    else:
        logger.error(f"Unsupported file type: {file_type}")
        return []

# Example usage
# pdf_text = load_document_to_text("path/to/your/file.pdf")
# csv_text = load_document_to_text("path/to/your/file.csv", csv_args={"delimiter": ","})
